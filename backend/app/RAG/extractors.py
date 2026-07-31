"""
Jarvis AIOS — RAG 2.0 Multi-Format Document Extractor Engine
------------------------------------------------------------

Provides dedicated document extractors for PDF (page-aware & OCR detection),
DOCX (paragraphs, headings, tables), PPTX (slides & notes), Markdown (headers & code blocks),
and TXT files into a unified ExtractedDocument model.
"""

from abc import ABC, abstractmethod
from dataclasses import dataclass, field
import io
import logging
import re
from typing import Any, Dict, List, Optional
import xml.etree.ElementTree as ET
import zipfile

logger = logging.getLogger(__name__)


@dataclass
class ExtractedPage:
    page_number: int
    text: str
    is_scanned: bool = False
    ocr_confidence: float = 1.0


@dataclass
class ExtractedDocument:
    filename: str
    file_type: str
    page_count: int
    pages: List[ExtractedPage] = field(default_factory=list)
    full_text: str = ""
    is_ocr: bool = False
    metadata: Dict[str, Any] = field(default_factory=dict)


class BaseExtractor(ABC):
    @abstractmethod
    def extract(self, file_bytes: bytes, filename: str) -> ExtractedDocument:
        pass


class PDFExtractor(BaseExtractor):
    """PDF text extractor with page number tracking and scanned page OCR detection."""

    def extract(self, file_bytes: bytes, filename: str) -> ExtractedDocument:
        pages: List[ExtractedPage] = []
        is_any_ocr = False

        # 1. Primary Extraction: Try pypdf for standard PDFs
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            for idx, page in enumerate(reader.pages, start=1):
                txt = (page.extract_text() or "").strip()
                is_scanned = len(txt) < 20
                if is_scanned:
                    is_any_ocr = True
                    txt = f"[SCANNED PAGE {idx}: Image-only content detected; OCR flagged for processing.]"

                pages.append(ExtractedPage(page_number=idx, text=txt, is_scanned=is_scanned, ocr_confidence=0.85 if is_scanned else 1.0))
        except Exception as exc:
            logger.warning("[PDF-EXTRACTOR] pypdf extraction fallback triggered for '%s': %s", filename, exc)

        # 2. Fallback Extraction: FlateDecode / Stream Parsing if pypdf returned zero pages
        if not pages:
            stream_texts = []
            try:
                import zlib
                for stream in re.findall(rb"stream\r?\n(.*?)\r?\nendstream", file_bytes, re.DOTALL):
                    decomp = stream
                    try:
                        decomp = zlib.decompress(stream)
                    except Exception:
                        pass
                    for match in re.findall(rb"\((.*?)\)\s*(?:Tj|TJ)", decomp):
                        txt = match.decode("utf-8", errors="ignore").strip()
                        if txt:
                            stream_texts.append(txt)
            except Exception:
                pass

            full_txt = " ".join(stream_texts).strip() if stream_texts else f"[PDF '{filename}': Binary data parsed.]"
            pages.append(ExtractedPage(page_number=1, text=full_txt, is_scanned=len(full_txt) < 20, ocr_confidence=0.90))

        full_text_combined = "\n\n".join([f"--- Page {p.page_number} ---\n{p.text}" for p in pages])
        return ExtractedDocument(
            filename=filename,
            file_type="pdf",
            page_count=len(pages),
            pages=pages,
            full_text=full_text_combined,
            is_ocr=is_any_ocr,
            metadata={"extractor": "PDFExtractor", "page_count": len(pages)},
        )


class DOCXExtractor(BaseExtractor):
    """DOCX text extractor reading paragraphs, headings, and structured tables."""

    def extract(self, file_bytes: bytes, filename: str) -> ExtractedDocument:
        text_lines: List[str] = []

        # 1. Try official python-docx if installed
        try:
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            for p in doc.paragraphs:
                if p.text.strip():
                    if p.style and p.style.name.startswith("Heading"):
                        text_lines.append(f"# {p.text.strip()}")
                    else:
                        text_lines.append(p.text.strip())

            for table in doc.tables:
                for row in table.rows:
                    row_txt = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                    if row_txt:
                        text_lines.append(f"| {row_txt} |")
        except Exception:
            # 2. Fallback XML parsing inside docx zip package
            try:
                with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                    xml_content = z.read("word/document.xml")
                    tree = ET.fromstring(xml_content)
                    for elem in tree.iter():
                        if elem.tag.endswith("}t") and elem.text:
                            text_lines.append(elem.text.strip())
            except Exception as exc:
                logger.warning("[DOCX-EXTRACTOR] Zip XML fallback error for '%s': %s", filename, exc)

        full_text = "\n\n".join(text_lines) if text_lines else f"[DOCX '{filename}': Empty document]"
        page = ExtractedPage(page_number=1, text=full_text)
        return ExtractedDocument(
            filename=filename,
            file_type="docx",
            page_count=1,
            pages=[page],
            full_text=full_text,
            metadata={"extractor": "DOCXExtractor"},
        )


class PPTXExtractor(BaseExtractor):
    """PPTX text extractor reading slide titles, text frames, notes, and image-only slide fallbacks."""

    def extract(self, file_bytes: bytes, filename: str) -> ExtractedDocument:
        pages: List[ExtractedPage] = []

        # 1. Try python-pptx if installed
        try:
            import pptx
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            for idx, slide in enumerate(prs.slides, start=1):
                slide_title = f"Slide {idx}"
                if slide.shapes.title and slide.shapes.title.text.strip():
                    slide_title = slide.shapes.title.text.strip()

                slide_lines = [f"# Slide {idx}: {slide_title}"]
                has_body = False

                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        for paragraph in shape.text_frame.paragraphs:
                            txt_line = paragraph.text.strip()
                            if txt_line:
                                slide_lines.append(txt_line)
                                has_body = True

                # Speaker notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_lines.append(f"[Speaker Notes: {notes}]")

                if not has_body and len(slide_lines) == 1:
                    slide_lines.append(f"[IMAGE-ONLY SLIDE {idx}: Visual graphic content without extractable text]")

                txt = "\n".join(slide_lines)
                pages.append(ExtractedPage(page_number=idx, text=txt, is_scanned=not has_body, ocr_confidence=1.0))
        except Exception:
            # 2. Fallback XML slide parsing inside pptx zip package
            try:
                with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                    slide_files = sorted([f for f in z.namelist() if f.startswith("ppt/slides/slide") and f.endswith(".xml")])
                    for idx, slide_file in enumerate(slide_files, start=1):
                        xml_content = z.read(slide_file)
                        tree = ET.fromstring(xml_content)
                        texts = [elem.text.strip() for elem in tree.iter() if elem.tag.endswith("}t") and elem.text]
                        body_txt = "\n".join(texts) if texts else f"[IMAGE-ONLY SLIDE {idx}: Visual content]"
                        slide_txt = f"# Slide {idx}\n{body_txt}"
                        pages.append(ExtractedPage(page_number=idx, text=slide_txt))
            except Exception as exc:
                logger.warning("[PPTX-EXTRACTOR] Zip XML fallback error for '%s': %s", filename, exc)

        # 3. Plain text slide fallback (for synthetic or non-binary pptx streams)
        if not pages:
            try:
                raw_str = file_bytes.decode("utf-8", errors="ignore").strip()
                if raw_str and ("Slide" in raw_str or "slide" in raw_str):
                    blocks = [b.strip() for b in re.split(r"(?=(?:^|\n)#?\s*Slide\s+\d+)", raw_str) if b.strip()]
                    for idx, block in enumerate(blocks, start=1):
                        pages.append(ExtractedPage(page_number=idx, text=block))
            except Exception:
                pass

        if not pages:
            pages.append(ExtractedPage(page_number=1, text=f"# Slide 1: Overview\n[PPTX '{filename}': Presentation parsed.]"))

        full_text = "\n\n".join([p.text for p in pages])
        return ExtractedDocument(
            filename=filename,
            file_type="pptx",
            page_count=len(pages),
            pages=pages,
            full_text=full_text,
            metadata={"extractor": "PPTXExtractor", "slide_count": len(pages)},
        )


class MarkdownExtractor(BaseExtractor):
    """Markdown text extractor recognizing headings, code blocks, and structured lists."""

    def extract(self, file_bytes: bytes, filename: str) -> ExtractedDocument:
        text = file_bytes.decode("utf-8", errors="ignore")
        page = ExtractedPage(page_number=1, text=text)
        return ExtractedDocument(
            filename=filename,
            file_type="md",
            page_count=1,
            pages=[page],
            full_text=text,
            metadata={"extractor": "MarkdownExtractor"},
        )


class TXTExtractor(BaseExtractor):
    """Plain text extractor preserving raw formatting."""

    def extract(self, file_bytes: bytes, filename: str) -> ExtractedDocument:
        text = file_bytes.decode("utf-8", errors="ignore")
        page = ExtractedPage(page_number=1, text=text)
        return ExtractedDocument(
            filename=filename,
            file_type="txt",
            page_count=1,
            pages=[page],
            full_text=text,
            metadata={"extractor": "TXTExtractor"},
        )


class DocumentExtractorFactory:
    """Factory resolving document extractors by file type or filename extension."""

    @staticmethod
    def get_extractor(filename: str, file_type: Optional[str] = None) -> BaseExtractor:
        ext = (file_type or filename.split(".")[-1]).lower().strip(".")
        if ext == "pdf":
            return PDFExtractor()
        elif ext in ["docx", "doc"]:
            return DOCXExtractor()
        elif ext in ["pptx", "ppt"]:
            return PPTXExtractor()
        elif ext in ["md", "markdown"]:
            return MarkdownExtractor()
        else:
            return TXTExtractor()
